import json
import os
import pathlib
import subprocess
from typing import Optional

import docx
import pptx
from pypdf import PdfReader

import anthropic
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
from fastapi.responses import HTMLResponse
from pydantic import BaseModel

load_dotenv()

app = FastAPI()
client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

INCLUDED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}
SYSTEM_DIRS = {"/System", "/Library", "/usr", "/bin", "/sbin", "/proc", "/dev", "/etc", "/var"}


def is_excluded(path: pathlib.Path) -> bool:
    if any(part.startswith(".") for part in path.parts[1:]):
        return True
    for sys_dir in SYSTEM_DIRS:
        if str(path).startswith(sys_dir):
            return True
    if path.suffix.lower() in {".app", ".exe", ".dmg", ".pkg"}:
        return True
    return False


def fuzzy_match(name: str, candidates: list) -> Optional[str]:
    name_l = name.lower().strip()
    for c in candidates:
        if c.lower() == name_l:
            return c
    for c in candidates:
        if c.lower().startswith(name_l):
            return c
    for c in candidates:
        if name_l in c.lower() or c.lower() in name_l:
            return c
    return None


def resolve_folder(parts: list, base: Optional[pathlib.Path] = None) -> Optional[pathlib.Path]:
    current = base or pathlib.Path.home()
    for part in parts:
        try:
            children = [p.name for p in current.iterdir() if p.is_dir()]
        except PermissionError:
            return None
        match = fuzzy_match(part, children)
        if not match:
            return None
        current = current / match
    return current


def collect_files(root: pathlib.Path) -> list:
    results = []
    for path in root.rglob("*"):
        if (
            path.is_file()
            and path.suffix.lower() in INCLUDED_EXTENSIONS
            and not is_excluded(path)
        ):
            results.append(str(path))
    return results


def extract_text(path: pathlib.Path, max_chars: int = 3000) -> str:
    """Extract a text snippet from a PDF, Word, or PowerPoint file."""
    try:
        ext = path.suffix.lower()
        if ext == ".pdf":
            reader = PdfReader(str(path))
            text = " ".join(page.extract_text() or "" for page in reader.pages)
        elif ext in {".docx", ".doc"}:
            doc = docx.Document(str(path))
            text = " ".join(p.text for p in doc.paragraphs)
        elif ext in {".pptx", ".ppt"}:
            prs = pptx.Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            text = " ".join(parts)
        else:
            return ""
        # Collapse whitespace and cap length
        text = " ".join(text.split())
        return text[:max_chars]
    except Exception:
        return ""


def parse_json_response(text: str):
    text = text.strip()
    if text.startswith("```"):
        text = text.split("```")[1]
        if text.startswith("json"):
            text = text[4:]
    return json.loads(text.strip())


class SearchRequest(BaseModel):
    query: str
    root_folders: list = []


class ComposeRequest(BaseModel):
    file_paths: list
    to_email: str
    subject: str = ""
    body: str = ""


@app.post("/search")
async def search(req: SearchRequest):
    parse_resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=300,
        messages=[
            {
                "role": "user",
                "content": (
                    'Parse this file search request and return JSON only.\n'
                    'Extract:\n'
                    '  "file_query": short description of the file they want\n'
                    '  "folder_parts": ordered list of folder name fragments in the path\n\n'
                    f'Request: "{req.query}"\n\n'
                    'Return only valid JSON, e.g. '
                    '{"file_query": "purchase agreement", "folder_parts": ["OneDrive", "Life and Home Buying"]}'
                ),
            }
        ],
    )

    try:
        parsed = parse_json_response(parse_resp.content[0].text)
    except (json.JSONDecodeError, IndexError):
        raise HTTPException(status_code=500, detail="Failed to parse query — try rephrasing.")

    file_query: str = parsed.get("file_query", "")
    folder_parts: list = parsed.get("folder_parts", [])

    roots = [pathlib.Path(r).expanduser() for r in req.root_folders if r.strip()] or [pathlib.Path.home()]
    folder_path = None
    for root in roots:
        candidate = resolve_folder(folder_parts, base=root) if folder_parts else root
        if candidate and candidate.exists():
            folder_path = candidate
            break

    if not folder_path:
        detail = f"Could not find folder: {' › '.join(folder_parts)}" if folder_parts else "Root folder not found."
        raise HTTPException(status_code=404, detail=detail)

    files = collect_files(folder_path)
    if not files:
        return {"matches": [], "folder": str(folder_path)}

    # --- Pass 1: filename-only ranking ---
    names_text = "\n".join(f for f in files)
    name_resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=800,
        messages=[
            {
                "role": "user",
                "content": (
                    f'I\'m looking for: "{file_query}"\n\n'
                    "From these file paths, return a JSON object with:\n"
                    '  "confident": true if 1+ filenames clearly match the query, false otherwise\n'
                    '  "matches": array of matching file paths (best first, max 8) — '
                    "only populated if confident is true\n\n"
                    "Return only valid JSON.\n\n"
                    + names_text
                ),
            }
        ],
    )

    try:
        name_result = parse_json_response(name_resp.content[0].text)
    except (json.JSONDecodeError, IndexError):
        name_result = {"confident": False, "matches": []}

    if name_result.get("confident") and name_result.get("matches"):
        ranked = [{"path": p, "excerpt": ""} for p in name_result["matches"]]
        return {"matches": ranked, "folder": str(folder_path), "searched": "filename"}

    # --- Pass 2: content-based ranking (no filename match found) ---
    candidates = files[:20]
    file_summaries = []
    for f in candidates:
        snippet = extract_text(pathlib.Path(f))
        file_summaries.append({
            "path": f,
            "name": pathlib.Path(f).name,
            "snippet": snippet or "(could not read content)",
        })

    summaries_text = "\n\n".join(
        f'FILE: {s["path"]}\nNAME: {s["name"]}\nCONTENT: {s["snippet"]}'
        for s in file_summaries
    )

    content_resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1500,
        messages=[
            {
                "role": "user",
                "content": (
                    f'I\'m looking for: "{file_query}"\n\n'
                    "Below are files with their content. Return a JSON array of objects for the "
                    "most relevant files (best match first, max 8). Only include files whose "
                    "content is relevant to the query. Each object must have:\n"
                    '  "path": the full file path\n'
                    '  "excerpt": a short 1-2 sentence quote or summary showing why it matches\n\n'
                    "Return only valid JSON.\n\n"
                    + summaries_text
                ),
            }
        ],
    )

    try:
        ranked = parse_json_response(content_resp.content[0].text)
    except (json.JSONDecodeError, IndexError):
        ranked = [{"path": f, "excerpt": ""} for f in files[:8]]

    return {"matches": ranked, "folder": str(folder_path), "searched": "content"}


@app.post("/compose")
async def compose(req: ComposeRequest):
    if not req.file_paths:
        raise HTTPException(status_code=400, detail="No files selected.")

    paths = [pathlib.Path(p) for p in req.file_paths]
    for p in paths:
        if not p.exists():
            raise HTTPException(status_code=404, detail=f"File not found: {p.name}")

    def esc(s: str) -> str:
        return s.replace("\\", "\\\\").replace('"', '\\"')

    first_name = paths[0].name
    subject = esc(req.subject or (first_name if len(paths) == 1 else f"{first_name} + {len(paths) - 1} more"))
    body = esc(req.body or "")
    to_email = esc(req.to_email)

    attachment_lines = "\n        ".join(
        f'make new attachment with properties {{file name:POSIX file "{esc(str(p))}"}} at after last paragraph'
        for p in paths
    )

    script = f'''
tell application "Mail"
    set newMsg to make new outgoing message with properties {{subject:"{subject}", content:"{body}", visible:true}}
    tell newMsg
        make new to recipient with properties {{address:"{to_email}"}}
        {attachment_lines}
    end tell
    activate
end tell
'''

    try:
        subprocess.run(["osascript", "-e", script], check=True, capture_output=True)
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"Could not open Mail: {e.stderr.decode()}")

    return {"success": True}


@app.get("/")
async def root():
    with open(pathlib.Path(__file__).parent / "index.html") as f:
        return HTMLResponse(f.read())
